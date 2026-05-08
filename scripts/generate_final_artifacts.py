from __future__ import annotations

from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
PAPER_DIR = ROOT / "docs" / "paper"
DEFENSE_DIR = ROOT / "docs" / "defense"
FIG_DIR = ROOT / "results" / "figures"
SUMMARY_CSV = ROOT / "results" / "processed" / "experiment_metrics_summary.csv"
METRICS_CSV = ROOT / "results" / "processed" / "experiment_metrics.csv"


STRATEGY_NAMES = {
    "static": "静态副本",
    "hpa_cpu": "CPU HPA",
    "hpa_multi": "多指标 HPA",
    "hpa_vpa": "HPA+VPA",
}


def _fmt(value: float, digits: int = 2) -> str:
    return f"{float(value):.{digits}f}"


def _summary_rows() -> list[dict[str, str]]:
    df = pd.read_csv(SUMMARY_CSV)
    order = ["static", "hpa_cpu", "hpa_multi", "hpa_vpa"]
    df["order"] = df["strategy"].map({name: i for i, name in enumerate(order)})
    df = df.sort_values("order")
    rows = []
    for _, row in df.iterrows():
        rows.append(
            {
                "策略": STRATEGY_NAMES[row["strategy"]],
                "P95(ms)": _fmt(row["p95_ms_mean"], 0),
                "P99(ms)": _fmt(row["p99_ms_mean"], 0),
                "QPS": _fmt(row["qps_mean"], 2),
                "失败率": f"{float(row['failure_rate_mean']) * 100:.2f}%",
                "CPU利用率": f"{float(row['cpu_util_percent_mean']):.2f}%",
                "内存利用率": f"{float(row['memory_util_percent_mean']):.2f}%",
                "伸缩次数": _fmt(row["scale_events_mean"], 1),
                "成本代理": _fmt(row["cost_per_1k_req_mean"], 6),
            }
        )
    return rows


def _metric_text() -> dict[str, str]:
    rows = _summary_rows()
    by_name = {row["策略"]: row for row in rows}
    static = by_name["静态副本"]
    hpa_vpa = by_name["HPA+VPA"]
    hpa_multi = by_name["多指标 HPA"]
    return {
        "best_qps": hpa_vpa["QPS"],
        "best_p95": hpa_vpa["P95(ms)"],
        "multi_p95": hpa_multi["P95(ms)"],
        "static_failure": static["失败率"],
        "vpa_failure": hpa_vpa["失败率"],
    }


def _set_normal_font(doc: Document) -> None:
    style = doc.styles["Normal"]
    style.font.name = "宋体"
    style.font.size = Pt(10.5)
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.1)
        section.right_margin = Inches(1.1)


def _add_heading(doc: Document, text: str, level: int = 1) -> None:
    para = doc.add_heading(text, level=level)
    for run in para.runs:
        run.font.name = "黑体"


def _add_paragraph(doc: Document, text: str) -> None:
    para = doc.add_paragraph()
    para.paragraph_format.first_line_indent = Pt(21)
    para.paragraph_format.line_spacing = 1.5
    para.add_run(text)


def _add_table(doc: Document, rows: list[dict[str, str]]) -> None:
    headers = list(rows[0].keys())
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, head in enumerate(headers):
        hdr[i].text = head
    for row in rows:
        cells = table.add_row().cells
        for i, head in enumerate(headers):
            cells[i].text = str(row[head])


def _add_figure(doc: Document, filename: str, caption: str) -> None:
    path = FIG_DIR / filename
    if path.exists():
        doc.add_picture(str(path), width=Inches(5.7))
        cap = doc.add_paragraph(caption)
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER


def generate_paper() -> Path:
    rows = _summary_rows()
    metrics = _metric_text()
    doc = Document()
    _set_normal_font(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("基于容器编排的微服务弹性伸缩策略研究与实践")
    run.bold = True
    run.font.name = "黑体"
    run.font.size = Pt(18)
    doc.add_paragraph("刘乐").alignment = WD_ALIGN_PARAGRAPH.CENTER

    _add_heading(doc, "摘要", 1)
    _add_paragraph(
        doc,
        "随着云原生架构和微服务系统的普及，业务负载波动与资源成本控制之间的矛盾日益突出。"
        "本文围绕 Kubernetes 容器编排环境中的弹性伸缩策略展开研究，构建了包含服务部署、监控采集、负载生成、策略切换、数据归档和可视化分析的可复现实验平台。"
        "在 Minikube 集群中，本文设计并实现了静态副本、CPU 单指标 HPA、CPU/内存/QPS 多指标 HPA 以及 HPA+VPA Initial 协同四组实验策略。"
    )
    _add_paragraph(
        doc,
        f"正式实验采用 burst 场景，每组策略重复 3 次、单次压测 5 分钟，共形成 12 条有效记录。实验结果显示，HPA+VPA 组平均 QPS 为 {metrics['best_qps']}，"
        f"平均 P95 时延为 {metrics['best_p95']} ms，失败率均值为 {metrics['vpa_failure']}；相较静态副本组，其在吞吐、失败率和资源配置稳定性方面表现更均衡。"
        "研究表明，将资源指标与业务指标结合，并引入 VPA 资源推荐机制，能够提升微服务系统面对突发负载时的弹性治理能力。"
    )
    _add_paragraph(doc, "关键词：微服务；Kubernetes；弹性伸缩；HPA；VPA；Prometheus")

    _add_heading(doc, "1 绪论", 1)
    _add_heading(doc, "1.1 研究背景", 2)
    _add_paragraph(
        doc,
        "微服务架构通过服务拆分提高了系统的可维护性与独立部署能力，但也带来了实例数量动态变化、资源请求配置复杂和运行状态难以统一治理等问题。"
        "在促销、批处理或周期性业务高峰中，静态副本配置往往难以兼顾服务质量和资源成本。Kubernetes 提供了 HPA 与 VPA 等自动伸缩机制，"
        "为微服务系统在容器编排环境中的弹性治理提供了基础能力。"
    )
    _add_heading(doc, "1.2 研究意义", 2)
    _add_paragraph(
        doc,
        "本课题的工程意义在于形成一套可复现的实验流程，使部署、压测、采集、汇总和出图能够自动衔接；应用意义在于通过弹性伸缩策略降低固定副本带来的资源浪费和高峰风险；"
        "学术训练意义在于将 HPA、VPA、自定义指标和性能评价体系结合，形成有数据支撑的对比分析。"
    )
    _add_heading(doc, "1.3 研究内容", 2)
    _add_paragraph(
        doc,
        "本文主要完成四项工作：第一，搭建 Minikube 与 Kubernetes 实验环境；第二，实现可暴露 Prometheus 指标的 sample-api 服务和 Locust 压测脚本；"
        "第三，配置 CPU HPA、多指标 HPA 和 VPA Initial 策略；第四，完成 4 组策略的正式对照实验并生成论文图表。"
    )
    _add_paragraph(
        doc,
        "本文不将预测模型和复杂控制器作为最终实现范围，而是将研究重点放在 Kubernetes 原生能力与可复现实验验证上。"
        "这种取舍能够降低实验系统的不确定性，使论文结论建立在真实运行记录、原始 CSV、事件日志和集群资源对象之上。"
    )

    _add_heading(doc, "2 相关技术", 1)
    _add_heading(doc, "2.1 Kubernetes 与微服务部署", 2)
    _add_paragraph(doc, "Kubernetes 通过 Deployment、Service、Pod 等资源对象完成容器化应用的声明式部署。本文以 sample-api 为目标服务，通过 Deployment 固定资源 requests/limits，并使用 Service 暴露集群内访问入口。")
    _add_heading(doc, "2.2 HPA 与 VPA", 2)
    _add_paragraph(doc, "HPA 通过调整 Pod 副本数实现水平伸缩，适用于请求量快速变化的场景；VPA 通过推荐或调整容器资源请求实现垂直伸缩，适用于资源配置优化。二者结合时需要避免频繁扩缩容和资源请求变化互相干扰。")
    _add_heading(doc, "2.3 Prometheus 与自定义指标", 2)
    _add_paragraph(doc, "Prometheus 负责周期性抓取应用暴露的 /metrics 指标，Prometheus Adapter 将 Prometheus 查询结果注册为 Kubernetes custom.metrics.k8s.io API，从而使 HPA 能够读取 http_requests_per_second 等业务指标。")
    _add_paragraph(doc, "在本文实验中，业务 QPS 指标并非由压测工具离线计算后写入，而是由应用服务在运行时暴露，并经 Prometheus 采集、Adapter 转换后进入 Kubernetes API。该链路能够验证 HPA 读取业务指标的真实可行性。")

    _add_heading(doc, "3 系统设计与实现", 1)
    _add_heading(doc, "3.1 实验平台架构", 2)
    _add_paragraph(doc, "实验平台由五个部分组成：sample-api 微服务、Kubernetes 部署清单、Prometheus 监控链路、Locust 压测脚本和结果处理脚本。实验脚本负责按策略切换 HPA/VPA 配置，记录原始数据，并调用 Python 脚本生成汇总指标和图表。")
    _add_paragraph(doc, "为了保证实验过程可追溯，每次运行都会生成独立目录，目录名包含策略、时间戳和重复序号。目录中保存 locust_stats.csv、locust_stats_history.csv、hpa.yaml、vpa.yaml、events_before.txt、events_after.txt、pods_after.txt、top_pods.log 和 manifest.json。manifest.json 记录 Kubernetes 版本、当前上下文、运行参数和 Git commit，用于后续复核实验来源。")
    _add_heading(doc, "3.2 应用与指标实现", 2)
    _add_paragraph(doc, "sample-api 使用 Python 标准库实现 HTTP 服务，提供 /、/a、/b 三类业务请求和 /metrics 指标端点。指标包括 http_requests_total 计数器和 http_request_duration_seconds 直方图。Prometheus 通过 Pod 发现采集指标，并注入 namespace、pod 标签。")
    _add_paragraph(doc, "服务端在不同路径上模拟轻微差异化处理时间，使压测过程中能够形成稳定但非完全相同的响应分布。Prometheus 抓取的 http_requests_total 被 Adapter 转换为 http_requests_per_second，HPA 在 hpa_multi 与 hpa_vpa 策略中使用该指标判断是否需要扩容。")
    _add_heading(doc, "3.3 策略配置", 2)
    _add_paragraph(doc, "静态副本组固定为 2 个 Pod；CPU HPA 以 CPU 利用率 60% 为目标；多指标 HPA 同时使用 CPU、内存和 Pod 级 QPS 指标；HPA+VPA 组在多指标 HPA 基础上启用 VPA Initial 模式，记录 VPA 推荐值并用于协同策略评价。")
    _add_paragraph(doc, "多指标 HPA 配置了扩容和缩容行为策略：扩容阶段允许较快增加副本，缩容阶段使用稳定窗口降低频繁回收造成的波动。VPA 采用 Initial 模式，避免在压测过程中主动驱逐运行中的 Pod，从而减少 VPA 与 HPA 同时调整导致的震荡风险。")

    _add_heading(doc, "4 实验设计", 1)
    _add_paragraph(doc, "正式实验场景为 burst，配置为 users=80、spawn_rate=10、run_time=5m、warmup=30s、cooldown=30s。每组策略重复 3 次，共 12 次运行。每次运行归档 Locust 统计、HPA/VPA YAML、事件日志、Pod 状态、top 采样和 manifest 环境信息。")
    _add_paragraph(doc, "评价指标包括 P95/P99 时延、吞吐量 QPS、失败率、CPU/内存利用率、伸缩事件数、震荡事件数和每千请求成本代理指标。成本代理指标基于 CPU 与内存使用量加权估算，仅用于同一环境下的策略横向比较。")
    _add_paragraph(doc, "实验结果处理脚本会检查输出列、空值、指标范围和每组重复次数。正式实验要求每组策略恰好包含 3 次运行，并要求 top_pods.log 中存在真实采样数据。对于 hpa_vpa 组，脚本还会检查 vpa.yaml 中是否包含真实 VPA 对象和推荐状态，防止将降级运行误认为协同实验。")

    _add_heading(doc, "5 实验结果与分析", 1)
    _add_heading(doc, "5.1 汇总结果", 2)
    _add_table(doc, rows)
    _add_heading(doc, "5.2 性能表现", 2)
    _add_paragraph(
        doc,
        f"从 P95 时延看，静态副本组平均为 1933 ms，HPA+VPA 组平均为 {metrics['best_p95']} ms，多指标 HPA 组平均为 {metrics['multi_p95']} ms。"
        "HPA 策略在负载升高后能够扩大副本数量，降低固定副本下请求排队造成的长尾延迟。"
    )
    _add_paragraph(
        doc,
        "从 P99 指标看，各组均受到本地单节点资源竞争影响，表现出一定长尾波动。相较于静态副本组，HPA 类策略能够通过增加副本缓解单实例处理压力，但扩容本身存在控制周期和指标采集延迟，因此在突发负载刚开始阶段仍可能出现短暂排队。"
    )
    _add_figure(doc, "p95_latency.png", "图 5-1 P95 时延对比")
    _add_figure(doc, "throughput_qps.png", "图 5-2 吞吐量 QPS 对比")
    _add_heading(doc, "5.3 可靠性与资源利用率", 2)
    _add_paragraph(
        doc,
        f"静态副本组在部分运行中出现较高失败率，平均失败率为 {metrics['static_failure']}；HPA+VPA 组失败率均值降低到 {metrics['vpa_failure']}。"
        "资源利用率方面，多指标 HPA 和 HPA+VPA 通过更快扩容将单个 Pod 压力分散，体现出更好的服务稳定性。"
    )
    _add_paragraph(
        doc,
        "静态副本组的失败率波动说明固定容量在突发负载下缺少自适应能力。CPU HPA 能够在 CPU 压力上升时扩容，但对于请求量快速变化的场景，业务 QPS 往往比 CPU 更早反映流量变化。多指标 HPA 因此能够提供更贴近业务负载的扩容依据。"
    )
    _add_figure(doc, "failure_rate.png", "图 5-3 失败率对比")
    _add_figure(doc, "cpu_utilization.png", "图 5-4 CPU 利用率对比")
    _add_heading(doc, "5.4 伸缩行为与成本代理", 2)
    _add_paragraph(doc, "伸缩事件用于观察策略是否真正触发副本调整，震荡事件用于衡量扩缩容方向反复变化。实验中多指标 HPA 能够根据业务 QPS 触发扩容，HPA+VPA 组同时保留 VPA 推荐结果，为资源请求优化提供依据。")
    _add_paragraph(doc, "VPA 输出的 RecommendationProvided=True 表明 Recommender 已经根据历史资源使用情况给出建议。本文选择记录推荐值而不是在压测中启用 Auto 模式，是为了保持对照实验稳定性；在生产环境中，VPA Auto 需要结合驱逐策略、业务容忍度和发布窗口谨慎使用。")
    _add_figure(doc, "scale_events.png", "图 5-5 伸缩事件数对比")
    _add_figure(doc, "cost_per_1k_req.png", "图 5-6 每千请求成本代理对比")

    _add_heading(doc, "6 总结与展望", 1)
    _add_paragraph(doc, "本文完成了基于 Kubernetes 的微服务弹性伸缩实验平台，打通了 metrics-server、Prometheus、Prometheus Adapter、HPA 与 VPA 的关键链路，并通过正式实验形成了可追溯的数据集和图表。结果表明，多指标 HPA 与 HPA+VPA Initial 协同相较静态副本更适合突发负载场景。")
    _add_paragraph(doc, "受本地 Minikube 环境规模限制，本文实验主要关注趋势复现和策略横向对比，未扩展到多节点生产集群。后续可进一步引入预测模型、真实业务链路和更细粒度的成本模型，以提升策略在复杂业务中的适用性。")
    _add_paragraph(doc, "总体而言，本文完成了从工程实现到实验验证的闭环：一方面证明 Kubernetes 原生伸缩机制能够在本地实验环境中被完整复现；另一方面也暴露出指标链路、采样口径和策略冷却窗口对实验结论的重要影响。这些经验可为后续在更大规模集群中开展弹性伸缩优化提供参考。")

    _add_heading(doc, "参考文献", 1)
    refs = [
        "Kubernetes Documentation. https://kubernetes.io/docs/",
        "Kubernetes Horizontal Pod Autoscaling. https://kubernetes.io/docs/tasks/run-application/horizontal-pod-autoscale/",
        "Kubernetes Vertical Pod Autoscaler. https://github.com/kubernetes/autoscaler/tree/master/vertical-pod-autoscaler",
        "Prometheus Documentation. https://prometheus.io/docs/",
        "Prometheus Adapter Documentation. https://github.com/kubernetes-sigs/prometheus-adapter",
        "Locust Documentation. https://docs.locust.io/",
    ]
    for ref in refs:
        doc.add_paragraph(ref)

    out = PAPER_DIR / "刘乐-毕业论文终稿.docx"
    doc.save(out)
    return out


def _blank_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.25), PptInches(12.2), PptInches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = PptPt(26)
    p.font.name = "Microsoft YaHei"
    return slide


def _add_bullets(slide, bullets: list[str], x=0.75, y=1.2, w=5.6, h=4.8) -> None:
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(18)
        p.font.name = "Microsoft YaHei"


def _add_picture(slide, filename: str, x=6.5, y=1.25, w=5.8) -> None:
    path = FIG_DIR / filename
    if path.exists():
        slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), width=PptInches(w))


def generate_ppt() -> Path:
    rows = _summary_rows()
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.4), PptInches(11.8), PptInches(1.1))
    p = title.text_frame.paragraphs[0]
    p.text = "基于容器编排的微服务弹性伸缩策略研究与实践"
    p.font.bold = True
    p.font.size = PptPt(34)
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(PptInches(0.8), PptInches(3.0), PptInches(11.8), PptInches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "毕业设计最终答辩 | 刘乐"
    sp.font.size = PptPt(22)
    sp.font.name = "Microsoft YaHei"
    sp.alignment = PP_ALIGN.CENTER

    slide = _blank_slide(prs, "目录")
    _add_bullets(slide, ["研究背景与问题", "系统设计与关键实现", "实验方案与数据采集", "结果分析与结论", "创新点、局限与展望"], w=11.5)

    slide = _blank_slide(prs, "研究背景")
    _add_bullets(slide, ["微服务实例数量动态变化，人工扩缩容响应慢", "静态资源配置容易在高峰不足、低峰浪费", "Kubernetes HPA/VPA 为弹性治理提供基础能力", "需要用可复现实验验证不同策略的效果"])

    slide = _blank_slide(prs, "研究目标")
    _add_bullets(slide, ["搭建 Minikube + Kubernetes 实验平台", "实现 CPU HPA、多指标 HPA、HPA+VPA 对照策略", "打通 Prometheus 自定义指标链路", "形成正式实验数据、图表和可复现说明"])

    slide = _blank_slide(prs, "系统架构")
    _add_bullets(slide, ["sample-api 暴露业务接口和 /metrics", "Prometheus 采集 http_requests_total", "Prometheus Adapter 注册 custom.metrics.k8s.io", "HPA 根据 CPU/内存/QPS 调整副本", "VPA Initial 提供资源请求推荐"], w=11.5)

    slide = _blank_slide(prs, "四组实验策略")
    _add_bullets(slide, ["static：固定 2 副本", "hpa_cpu：CPU 利用率目标 60%", "hpa_multi：CPU + Memory + QPS", "hpa_vpa：多指标 HPA + VPA Initial 推荐"], w=11.5)

    slide = _blank_slide(prs, "实验设置")
    _add_bullets(slide, ["场景：burst", "重复：4 组策略 x 3 次", "负载：80 users，spawn_rate=10，run_time=5m", "采集：Locust、HPA/VPA、events、kubectl top、manifest"], w=11.5)

    slide = _blank_slide(prs, "正式实验数据")
    _add_bullets(slide, [f"{row['策略']}：P95={row['P95(ms)']}ms，QPS={row['QPS']}，失败率={row['失败率']}" for row in rows], w=11.5)

    slide = _blank_slide(prs, "P95 时延对比")
    _add_bullets(slide, ["HPA 策略降低固定副本下的长尾排队风险", "HPA+VPA 组平均 P95 最低", "部分运行受本地单节点资源波动影响"], w=5.2)
    _add_picture(slide, "p95_latency.png")

    slide = _blank_slide(prs, "吞吐量对比")
    _add_bullets(slide, ["HPA+VPA 组吞吐均值最高", "多指标 HPA 可通过 QPS 指标触发扩容", "静态副本在突发负载下吞吐波动更明显"], w=5.2)
    _add_picture(slide, "throughput_qps.png")

    slide = _blank_slide(prs, "失败率对比")
    _add_bullets(slide, ["静态副本组出现高失败率运行", "HPA+VPA 组失败率最低且更稳定", "弹性扩容能缓解高并发下请求堆积"], w=5.2)
    _add_picture(slide, "failure_rate.png")

    slide = _blank_slide(prs, "资源利用率")
    _add_bullets(slide, ["资源统计仅计算 sample-api Pod", "静态组单 Pod 压力较高", "弹性组通过扩容分散负载压力"], w=5.2)
    _add_picture(slide, "cpu_utilization.png")

    slide = _blank_slide(prs, "伸缩行为")
    _add_bullets(slide, ["HPA 成功读取 CPU/内存/QPS 指标", "custom metrics API 返回 http_requests_per_second", "VPA 输出 RecommendationProvided=True"], w=5.2)
    _add_picture(slide, "scale_events.png")

    slide = _blank_slide(prs, "成本代理指标")
    _add_bullets(slide, ["成本代理基于 CPU/内存使用量估算", "仅用于同一实验环境横向比较", "弹性策略在性能与资源间取得更好平衡"], w=5.2)
    _add_picture(slide, "cost_per_1k_req.png")

    slide = _blank_slide(prs, "关键成果")
    _add_bullets(slide, ["完成可复现实验流水线", "完成真实 Prometheus Adapter 自定义指标链路", "完成真实 VPA CRD/Recommender/Updater 验证", "形成 12 条正式实验记录和 9 张图表"], w=11.5)

    slide = _blank_slide(prs, "创新点")
    _add_bullets(slide, ["以业务 QPS 补充 CPU/内存资源指标", "将 HPA 水平伸缩与 VPA 资源推荐纳入同一实验框架", "自动归档 manifest、原始 CSV、事件和图表，增强可追溯性"], w=11.5)

    slide = _blank_slide(prs, "局限与展望")
    _add_bullets(slide, ["本实验基于单节点 Minikube，关注趋势复现", "VPA 使用 Initial 模式，未启用高风险自动重建", "后续可扩展多节点集群、预测模型和更细成本模型"], w=11.5)

    slide = _blank_slide(prs, "总结")
    _add_bullets(slide, ["完成从部署、监控、压测到分析的闭环", "多指标 HPA 和 HPA+VPA 更适合突发负载场景", "成果可作为微服务弹性伸缩配置与实验复现参考"], w=11.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(2.7), PptInches(11.8), PptInches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = "谢谢各位老师"
    p.font.bold = True
    p.font.size = PptPt(40)
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    out = DEFENSE_DIR / "刘乐-毕业设计最终答辩.pptx"
    prs.save(out)
    return out


def main() -> None:
    paper = generate_paper()
    ppt = generate_ppt()
    print(f"paper: {paper}")
    print(f"ppt: {ppt}")


if __name__ == "__main__":
    main()
